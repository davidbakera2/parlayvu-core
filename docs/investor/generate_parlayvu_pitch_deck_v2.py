from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUTPUT = OUT_DIR / "ParlayVU_Angel_Pitch_Deck_v2.pptx"

W, H = 13.333, 7.5

COLORS = {
    "ink": RGBColor(20, 31, 28),
    "forest": RGBColor(30, 53, 47),
    "evergreen": RGBColor(42, 74, 66),
    "sage": RGBColor(124, 155, 143),
    "mist": RGBColor(226, 235, 229),
    "cream": RGBColor(248, 244, 236),
    "white": RGBColor(255, 255, 255),
    "copper": RGBColor(169, 94, 50),
    "clay": RGBColor(211, 143, 87),
    "slate": RGBColor(74, 91, 86),
    "line": RGBColor(206, 219, 211),
    "soft": RGBColor(239, 234, 224),
}

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def rgb(name: str) -> RGBColor:
    return COLORS[name]


def add_text(slide, x, y, w, h, text, size=16, color="ink", bold=False, align=None, font=FONT_BODY):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return shape


def add_kicker(slide, text, dark=False):
    color = "clay" if dark else "copper"
    add_text(slide, 0.72, 0.5, 4.0, 0.22, text.upper(), 8.5, color, True)


def add_title(slide, kicker, title, subtitle=None, dark=False):
    add_kicker(slide, kicker, dark)
    add_text(slide, 0.7, 0.86, 7.2, 0.95, title, 27, "white" if dark else "forest", True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, 0.73, 1.78, 6.9, 0.5, subtitle, 12.5, "mist" if dark else "slate")


def add_bg(slide, dark=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb("forest" if dark else "cream")
    if dark:
        orb1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.7), Inches(-0.9), Inches(4.6), Inches(4.6))
        orb1.fill.solid()
        orb1.fill.fore_color.rgb = rgb("evergreen")
        orb1.fill.transparency = 12
        orb1.line.fill.background()
        orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.1), Inches(4.5), Inches(2.8), Inches(2.8))
        orb2.fill.solid()
        orb2.fill.fore_color.rgb = rgb("copper")
        orb2.fill.transparency = 62
        orb2.line.fill.background()
    else:
        wash = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.7), Inches(-0.4), Inches(3.2), Inches(3.2))
        wash.fill.solid()
        wash.fill.fore_color.rgb = rgb("mist")
        wash.fill.transparency = 8
        wash.line.fill.background()


def add_wordmark(slide, dark=False):
    color = "white" if dark else "forest"
    add_text(slide, 10.75, 0.43, 1.85, 0.25, "ParlayVU.ai", 11, color, True, PP_ALIGN.RIGHT, FONT_HEAD)


def add_footer(slide, n, total, dark=False):
    color = "mist" if dark else "slate"
    add_text(slide, 0.72, 7.04, 2.6, 0.18, "Agentic client-work operating system", 7.5, color)
    add_text(slide, 11.72, 7.04, 0.85, 0.18, f"{n:02d}/{total:02d}", 7.5, color, False, PP_ALIGN.RIGHT)


def rect(slide, x, y, w, h, fill="white", line="line", radius=True, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.8)
    return shape


def card(slide, x, y, w, h, label, body, fill="white", accent="copper", dark=False):
    rect(slide, x, y, w, h, fill, "line")
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.24, label.upper(), 7.3, accent, True)
    add_text(slide, x + 0.18, y + 0.47, w - 0.36, h - 0.56, body, 11.3, "white" if dark else "ink")


def pill(slide, x, y, w, h, text, fill="forest", color="white", size=10.3):
    shape = rect(slide, x, y, w, h, fill, fill)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = rgb(color)
    return shape


def shape_text(shape, text, size=12, color="ink", bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT_BODY
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)


def arrow(slide, x, y, w, h, text, fill="mist", text_color="forest"):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    sh.line.fill.background()
    shape_text(sh, text, 10.5, text_color, True)
    return sh


def circle(slide, x, y, d, text, fill="forest", color="white", size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    sh.line.color.rgb = rgb("line")
    shape_text(sh, text, size, color, True)
    return sh


def line(slide, x1, y1, x2, y2, color="line", width=1.4):
    sh = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    sh.line.color.rgb = rgb(color)
    sh.line.width = Pt(width)
    return sh


def slide_title(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    add_wordmark(slide, True)
    add_kicker(slide, "Angel / Pre-Seed", True)
    add_text(slide, 0.72, 0.95, 7.1, 1.35, "AI digital workers for client work", 35, "white", True, font=FONT_HEAD)
    add_text(
        slide,
        0.76,
        2.35,
        6.45,
        0.75,
        "ParlayVU turns meetings, source material, approvals, and follow-up into an operating system for expert firms.",
        15,
        "mist",
    )
    pill(slide, 0.78, 3.45, 2.05, 0.38, "$100k-$300k raise", "copper")
    pill(slide, 2.98, 3.45, 2.85, 0.38, "Teams/M365 beachhead", "evergreen")
    pill(slide, 5.98, 3.45, 2.25, 0.38, "RamAir pilot proof", "sage", "ink")
    for x, y, label in [(8.5, 1.05, "Brief"), (10.25, 1.82, "Nathan"), (8.9, 3.05, "Agents"), (10.55, 4.0, "Approvals"), (8.05, 4.9, "Outputs")]:
        circle(slide, x, y, 1.05, label, "cream" if label == "Nathan" else "evergreen", "forest" if label == "Nathan" else "white", 9.5)
    line(slide, 9.52, 1.58, 10.22, 2.16, "sage")
    line(slide, 10.42, 2.86, 9.35, 3.28, "sage")
    line(slide, 9.85, 3.63, 10.6, 4.2, "sage")
    line(slide, 10.45, 4.78, 9.0, 5.15, "sage")
    add_text(slide, 8.35, 6.15, 3.75, 0.32, "VU is your brand. Parlay your VU.", 14, "white", True, PP_ALIGN.CENTER, FONT_HEAD)
    add_footer(slide, idx, total, True)


def slide_problem(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Problem", "Expert firms lose margin in the handoffs", "The expensive work begins after the client call ends.")
    circle(slide, 5.8, 3.0, 1.15, "Client\ncall", "forest", "white", 12)
    nodes = [
        (1.0, 2.05, "Notes", "transcripts, recaps"),
        (1.35, 4.55, "Tasks", "owners, deadlines"),
        (4.0, 5.55, "Approvals", "claims, sends, deploys"),
        (8.25, 5.2, "Reporting", "dashboards, updates"),
        (10.0, 3.05, "Content", "posts, pages, emails"),
        (8.65, 1.45, "Follow-up", "client-ready next steps"),
        (3.55, 1.2, "Files", "source material"),
    ]
    for x, y, label, detail in nodes:
        card(slide, x, y, 2.25, 0.88, label, detail, "white")
        line(slide, 6.35, 3.55, x + 1.12, y + 0.44, "sage", 1.1)
    add_text(slide, 0.85, 6.25, 10.2, 0.34, "Generic AI drafts one asset. The margin leak is coordinating the whole client workflow.", 13.5, "forest", True)
    add_footer(slide, idx, total)


def slide_positioning(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Positioning", "Not Microsoft-dependent. Microsoft-native first.", "ParlayVU is an operating layer for client work; Teams and M365 are the launch surface customers already use.")
    add_text(slide, 0.95, 2.65, 4.15, 0.75, "Category thesis", 10, "copper", True)
    add_text(slide, 0.95, 3.05, 4.6, 1.05, "AI digital workers need memory, workflow, and approval context before they can run real client work.", 21, "forest", True, font=FONT_HEAD)
    items = [
        ("Beachhead", "Teams, SharePoint Files, Word, Planner direction, Power BI starters"),
        ("Core OS", "Nathan orchestration, project memory, specialist agents, approval gates"),
        ("Expansion", "Client portals, websites, email, analytics, publishing, partner channels"),
    ]
    y = 2.55
    for label, body in items:
        card(slide, 6.45, y, 5.45, 0.98, label, body, "white")
        y += 1.18
    pill(slide, 6.45, 6.1, 5.45, 0.42, "Keep Microsoft as the wedge, not the whole company story.", "forest")
    add_footer(slide, idx, total)


def slide_workflow(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Workflow", "From raw interaction to approved client work", "Shorter path from source material to governed outputs.")
    steps = [
        ("Brief\nIntake", "audience, offer, constraints"),
        ("Nathan\nOrchestrates", "routes work, checks memory"),
        ("Specialist\nAgents", "write, design, build, analyze"),
        ("Approval\nGates", "human control on risk"),
        ("Client-Ready\nOutputs", "notes, assets, sites, reports"),
    ]
    x = 0.8
    for i, (label, detail) in enumerate(steps):
        arrow(slide, x, 3.0, 2.18, 0.88, label, "forest" if i == 1 else "mist", "white" if i == 1 else "forest")
        add_text(slide, x + 0.08, 4.05, 1.8, 0.45, detail, 8.8, "slate", False, PP_ALIGN.CENTER)
        x += 2.35
    rect(slide, 1.35, 5.18, 10.4, 0.65, "white", "line")
    add_text(slide, 1.65, 5.39, 9.8, 0.2, "Project memory learns from source material, decisions, approvals, outputs, and results.", 12.2, "forest", True, PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_architecture(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Product Architecture", "A client-work OS layered on existing tools", "Editable schematic, not a claim of every integration being production-complete today.")
    layers = [
        ("Front doors", "Teams messages • M365 Files • email/task surfaces • web/demo UI", "mist"),
        ("Orchestration", "Nathan • FastAPI/LangGraph • specialist agent registry • readiness checks", "white"),
        ("Memory + governance", "clients • projects • source assets • approval states • agent events", "mist"),
        ("Output systems", "meeting notes • approval packets • campaign sites • dashboard starters • follow-up", "white"),
    ]
    y = 2.35
    for label, body, fill in layers:
        rect(slide, 1.0, y, 11.05, 0.76, fill, "line")
        add_text(slide, 1.28, y + 0.18, 2.0, 0.25, label.upper(), 8, "copper", True)
        add_text(slide, 3.2, y + 0.18, 8.3, 0.28, body, 11.2, "ink")
        y += 0.93
    pill(slide, 1.0, 6.25, 5.25, 0.38, "Working proof: routing, memory scaffold, Files outputs, approvals", "forest", "white", 9)
    pill(slide, 6.55, 6.25, 5.5, 0.38, "Roadmap: production native Teams audio/video media bridge", "soft", "forest", 9)
    add_footer(slide, idx, total)


def slide_ramair(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Pilot Proof", "RamAir shows the repeatable content engine", "One weekly podcast becomes a coordinated client workspace and distribution cadence.")
    circle(slide, 0.95, 3.0, 1.15, "Podcast\nsource", "forest", "white", 10.5)
    outputs = [
        ("YouTube + site", 3.0, 1.75),
        ("Short clips", 5.15, 1.75),
        ("Social posts", 7.3, 1.75),
        ("Transcript insights", 3.0, 3.35),
        ("Case studies", 5.15, 3.35),
        ("Sell sheets", 7.3, 3.35),
    ]
    for label, x, y in outputs:
        card(slide, x, y, 1.75, 0.95, "Output", label, "white")
        line(slide, 2.05, 3.58, x, y + 0.47, "sage", 1.0)
    rect(slide, 9.75, 1.62, 2.35, 3.8, "forest", "forest")
    add_text(slide, 10.02, 1.95, 1.75, 0.28, "CLIENT SIGNAL", 8, "clay", True, PP_ALIGN.CENTER)
    add_text(slide, 10.0, 2.38, 1.85, 1.7, "“One show gives us a week’s worth of quality content across site and socials.”", 15, "white", True, PP_ALIGN.CENTER, FONT_HEAD)
    add_text(slide, 9.98, 4.42, 1.9, 0.45, "Position as proof of workflow pattern, not broad SaaS traction.", 8.3, "mist", False, PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_proof_boundaries(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Credible Proof", "Be specific about what works now", "Investor confidence improves when the roadmap is clearly separated from current proof.")
    add_text(slide, 1.05, 2.35, 4.3, 0.35, "Working / demo-ready", 15, "forest", True, font=FONT_HEAD)
    working = [
        "Nathan routing and specialist agent registry",
        "Project-memory scaffold for client work",
        "M365 Files meeting notes as .md and .docx",
        "Approval endpoints and Teams-ready cards",
        "RamAir workspace artifacts and dashboard starter",
        "Provider-hosted Tavus/Nathan validation",
    ]
    y = 2.95
    for item in working:
        pill(slide, 1.05, y, 4.8, 0.32, item, "white", "forest", 8.6)
        y += 0.45
    add_text(slide, 7.0, 2.35, 4.3, 0.35, "Roadmap / do not overclaim", 15, "forest", True, font=FONT_HEAD)
    roadmap = [
        "Autonomous native Teams A/V participation",
        "Graph media injection into live calls",
        "Fully automated outbound sends",
        "Mature SaaS traction before paid pilots or LOIs",
        "Fully automated reporting refresh before wiring",
    ]
    y = 2.95
    for item in roadmap:
        pill(slide, 7.0, y, 4.95, 0.32, item, "soft", "ink", 8.6)
        y += 0.45
    add_footer(slide, idx, total)


def slide_market(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Market Wedge", "Start where expertise becomes deliverables", "A narrow wedge with repeatable workflows and urgent margin pressure.")
    circle(slide, 5.55, 2.65, 2.25, "Ideal\npilot", "forest", "white", 18)
    rings = [
        ("Expert-led firms", 1.05, 2.2),
        ("M365 / Teams users", 3.25, 1.45),
        ("Repeatable meetings", 8.25, 1.45),
        ("Content + reporting", 10.15, 2.2),
        ("Approval-heavy work", 2.1, 4.9),
        ("Founder-led sales fit", 9.0, 4.9),
    ]
    for label, x, y in rings:
        pill(slide, x, y, 2.3, 0.45, label, "white", "forest", 9.2)
    add_text(slide, 1.05, 6.1, 10.75, 0.36, "Michigan/Ann Arbor is the launch network and credibility base, not the market ceiling.", 13.5, "forest", True, PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_business_model(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Business Model", "High-touch pilots reveal subscription modules", "Early services help discover repeatable product packaging.")
    tiers = [
        ("1", "Pilot setup", "configure workspace, memory, approvals, demo workflow"),
        ("2", "Monthly workspace", "active clients/projects, agent seats, integrations"),
        ("3", "Premium modules", "publishing, dashboards, site generation, advanced memory"),
    ]
    x = 1.05
    for num, label, body in tiers:
        rect(slide, x, 2.65, 3.25, 2.45, "white", "line")
        circle(slide, x + 1.1, 2.0, 0.75, num, "forest", "white", 16)
        add_text(slide, x + 0.35, 3.0, 2.55, 0.35, label, 18, "forest", True, PP_ALIGN.CENTER, FONT_HEAD)
        add_text(slide, x + 0.42, 3.68, 2.4, 0.75, body, 11.5, "slate", False, PP_ALIGN.CENTER)
        x += 3.75
    add_text(slide, 1.2, 5.95, 10.6, 0.35, "The long-term product is a repeatable client-workspace package for agencies and expert firms.", 13.3, "forest", True, PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_gtm(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Go To Market", "Founder-led pilots create the first playbook", "Sell one concrete workflow before expanding the platform story.")
    steps = [
        ("Introductions", "angels, operators, agencies"),
        ("RamAir-style pilot", "one workflow, visible outputs"),
        ("Proof metrics", "time saved, output cadence, approvals"),
        ("Templates", "onboarding, notes, packets, dashboards"),
        ("Referrals", "same buyer pain, adjacent firms"),
    ]
    x = 0.8
    for label, body in steps:
        card(slide, x, 2.85, 2.1, 1.42, label, body, "white")
        x += 2.38
    add_text(slide, 1.1, 5.25, 10.7, 0.7, "Pilot promise: turn client meetings and source material into approved deliverables and follow-up inside the tools the firm already uses.", 18, "forest", True, PP_ALIGN.CENTER, FONT_HEAD)
    add_footer(slide, idx, total)


def slide_differentiation(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Differentiation", "Memory + governance + outputs in one workflow", "ParlayVU is not trying to win as another blank chat box.")
    add_text(slide, 1.0, 2.35, 4.6, 0.25, "Point tools", 10, "copper", True, PP_ALIGN.CENTER)
    add_text(slide, 7.25, 2.35, 4.6, 0.25, "Operating system", 10, "copper", True, PP_ALIGN.CENTER)
    labels = [
        (1.0, 2.8, "Generic AI", "drafts content, loses workflow context"),
        (1.0, 4.45, "Content tools", "help make assets, not client operations"),
        (7.25, 2.8, "ParlayVU", "coordinates client work across memory, agents, approvals, and outputs"),
        (7.25, 4.45, "Automation tools", "connect systems, but still need agentic judgment and governance"),
    ]
    for x, y, label, body in labels:
        fill = "forest" if label == "ParlayVU" else "white"
        card(slide, x, y, 4.6, 1.15, label, body, fill, "clay" if label == "ParlayVU" else "copper", label == "ParlayVU")
    line(slide, 6.45, 2.6, 6.45, 5.85, "sage", 2)
    add_footer(slide, idx, total)


def slide_roadmap(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Roadmap", "90 days to pilot-ready product", "Keep the media-bot ambition credible by sequencing it after core workflow hardening.")
    milestones = [
        ("30 days", "Polish RamAir workspace, meeting notes, approvals, readiness, demo script"),
        ("60 days", "Package onboarding templates and run 3-5 qualified pilot conversations"),
        ("90 days", "Paid pilots or LOIs, measurable outputs, pricing and module clarity"),
        ("Later", "Validate native Teams roster/media bridge after Graph and provider terms are proven"),
    ]
    x = 0.95
    for label, body in milestones:
        circle(slide, x + 0.55, 3.05, 0.88, label, "forest" if label != "Later" else "soft", "white" if label != "Later" else "forest", 9.5)
        card(slide, x, 4.1, 2.55, 1.18, "Milestone", body, "white")
        x += 3.0
    line(slide, 1.9, 3.5, 10.95, 3.5, "sage", 2.2)
    add_footer(slide, idx, total)


def slide_ask(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_wordmark(slide)
    add_title(slide, "Ask", "Raising $100k-$300k to turn proof into revenue", "Capital plus pilot/customer introductions are equally valuable at this stage.")
    circle(slide, 1.05, 2.55, 2.3, "$100k\n-\n$300k", "forest", "white", 18)
    uses = [
        ("Product hardening", 4.1, 2.25),
        ("Pilot delivery", 7.05, 2.25),
        ("M365 / Teams integration", 9.9, 2.25),
        ("Avatar validation", 4.1, 4.2),
        ("Pitch + design assets", 7.05, 4.2),
        ("Founder-led GTM", 9.9, 4.2),
    ]
    for label, x, y in uses:
        pill(slide, x, y, 2.35, 0.52, label, "white", "forest", 9.5)
    add_text(slide, 1.05, 5.95, 10.95, 0.42, "Target milestone: RamAir polished, 3-5 qualified pilots active, first paid pilots or signed LOIs within 90 days of funding.", 13, "forest", True, PP_ALIGN.CENTER)
    add_footer(slide, idx, total)


def slide_close(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    add_wordmark(slide, True)
    add_text(slide, 1.1, 1.35, 8.5, 1.05, "VU is your brand.\nParlay your VU.", 39, "white", True, font=FONT_HEAD)
    add_text(slide, 1.15, 3.2, 5.9, 0.75, "ParlayVU turns every client interaction into approved work, reusable memory, and distribution-ready output.", 18, "mist")
    rect(slide, 7.8, 2.0, 3.8, 2.65, "cream", "cream")
    add_text(slide, 8.15, 2.38, 3.05, 0.35, "Investor asks", 20, "forest", True, PP_ALIGN.CENTER, FONT_HEAD)
    add_text(slide, 8.25, 3.08, 2.8, 0.85, "Small checks\nPilot introductions\nMicrosoft 365 operator/customer intros", 14, "ink", False, PP_ALIGN.CENTER)
    add_footer(slide, idx, total, True)


BUILDERS = [
    slide_title,
    slide_problem,
    slide_positioning,
    slide_workflow,
    slide_architecture,
    slide_ramair,
    slide_proof_boundaries,
    slide_market,
    slide_business_model,
    slide_gtm,
    slide_differentiation,
    slide_roadmap,
    slide_ask,
    slide_close,
]


def validate(path: Path, expected: int) -> None:
    with ZipFile(path) as zf:
        slide_xml = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
    if len(slide_xml) != expected:
        raise RuntimeError(f"Expected {expected} slide XML files, found {len(slide_xml)}")
    prs = Presentation(str(path))
    if len(prs.slides) != expected:
        raise RuntimeError(f"python-pptx reload found {len(prs.slides)} slides")
    if not any(name.startswith("ppt/presentation.xml") for name in zf.namelist()):
        raise RuntimeError("Missing presentation.xml")


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    total = len(BUILDERS)
    for i, builder in enumerate(BUILDERS, start=1):
        builder(prs, i, total)
    prs.save(OUTPUT)
    validate(OUTPUT, total)
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out}")
    print(f"Slides: {len(BUILDERS)}")
