from __future__ import annotations

from pathlib import Path
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "parlayvu_site_assets"
OUTPUT = OUT_DIR / "ParlayVU_Angel_Pitch_Deck_v5.pptx"

W, H = 13.333, 7.5
# parlayvu.ai hero / body background (global.css)
HERO_BG_DEEP = RGBColor(3, 9, 18)  # --bg-deep #030912
HERO_BG = RGBColor(6, 17, 31)  # --bg #06111f
HERO_BG_END = RGBColor(7, 21, 40)  # linear-gradient end #071528
GRID_STEP_IN = 56 / 96  # site grid: 56px at 96dpi

COLORS = {
    "night": RGBColor(3, 9, 18),
    "navy": RGBColor(6, 17, 31),
    "panel": RGBColor(7, 21, 40),
    "panel2": RGBColor(12, 34, 64),
    "aqua": RGBColor(72, 224, 208),
    "blue": RGBColor(105, 167, 255),
    "violet": RGBColor(157, 123, 255),
    "lime": RGBColor(199, 243, 109),
    "white": RGBColor(247, 251, 255),
    "muted": RGBColor(199, 213, 232),
    "soft": RGBColor(159, 179, 203),
    "line": RGBColor(180, 210, 255),
    "ink": RGBColor(2, 16, 28),
    "green": RGBColor(32, 53, 47),
    "orange": RGBColor(217, 144, 79),
}

FONT_HEAD = "Inter"
FONT_BODY = "Inter"

ASSETS = {
    "logo": ("https://parlayvu.ai/brand/parlayvu-logo-dark.png", "parlayvu-logo-dark.png"),
    "icon": ("https://parlayvu.ai/brand/parlayvu-site-icon.png", "parlayvu-site-icon.png"),
    "nathan": ("https://parlayvu.ai/agents/Nathan.png", "Nathan.png"),
    "ramair": ("https://parlayvu.ai/case-studies/ramair-david-hart.png", "ramair-david-hart.png"),
    "alex": ("https://parlayvu.ai/agents/Alex.png", "Alex.png"),
    "ava": ("https://parlayvu.ai/agents/Ava.png", "Ava.png"),
    "dylan": ("https://parlayvu.ai/agents/Dylan.png", "Dylan.png"),
}


def rgb(name: str) -> RGBColor:
    return COLORS[name]


def ensure_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(exist_ok=True)
    paths: dict[str, Path] = {}
    for key, (url, filename) in ASSETS.items():
        path = ASSET_DIR / filename
        if not path.exists():
            try:
                request = Request(
                    url,
                    headers={
                        "User-Agent": (
                            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                            "AppleWebKit/537.36 (KHTML, like Gecko) "
                            "Chrome/125.0 Safari/537.36"
                        )
                    },
                )
                with urlopen(request, timeout=20) as response:
                    path.write_bytes(response.read())
            except (OSError, HTTPError, URLError):
                continue
        paths[key] = path
    return paths


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=16,
    color="white",
    bold=False,
    align=None,
    font=FONT_BODY,
    anchor=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return shape


def rect(slide, x, y, w, h, fill="panel", line="line", radius=True, transparency=0, line_transparency=72):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line)
    shape.line.transparency = line_transparency
    shape.line.width = Pt(0.85)
    return shape


def pill(slide, x, y, w, h, text, fill="aqua", color="ink", size=8.7, bold=True, transparency=0):
    shape = rect(slide, x, y, w, h, fill, fill, True, transparency, 100)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return shape


def _hero_glow(slide, cx_pct, cy_pct, radius_pct, color: RGBColor, alpha: float):
    """Radial glow matching parlayvu.ai body background gradients."""
    d = radius_pct * max(W, H) * 2
    x = cx_pct * W - d / 2
    y = cy_pct * H - d / 2
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = max(0, min(100, int((1 - alpha) * 100)))
    shape.line.fill.background()


def _hero_grid(slide):
    """Faint 56px grid with top-weighted fade (parlay-shell::before)."""
    grid_rgb = RGBColor(255, 255, 255)
    fade_end = H * 0.78
    x = 0.0
    while x <= W:
        line = slide.shapes.add_connector(1, Inches(x), Inches(0), Inches(x), Inches(fade_end))
        line.line.color.rgb = grid_rgb
        line.line.width = Pt(0.45)
        line.line.transparency = 96
        x += GRID_STEP_IN
    y = 0.0
    while y <= fade_end:
        fade = y / fade_end if fade_end else 0
        transparency = 96 + int(fade * 3)
        line = slide.shapes.add_connector(1, Inches(0), Inches(y), Inches(W), Inches(y))
        line.line.color.rgb = grid_rgb
        line.line.width = Pt(0.45)
        line.line.transparency = min(100, transparency)
        y += GRID_STEP_IN


def add_bg(slide, variant="dark"):
    """parlayvu.ai hero-section background: gradient base + radial glows + masked grid."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = HERO_BG_DEEP

    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(H))
    base.line.fill.background()
    base.fill.gradient()
    base.fill.gradient_angle = 90
    base.fill.gradient_stops[0].position = 0.0
    base.fill.gradient_stops[0].color.rgb = HERO_BG_DEEP
    base.fill.gradient_stops[1].position = 1.0
    base.fill.gradient_stops[1].color.rgb = HERO_BG_END

    # body { radial-gradient(...) } — same positions as global.css
    _hero_glow(slide, 0.14, 0.10, 0.28, rgb("aqua"), 0.16)
    _hero_glow(slide, 0.86, 0.08, 0.30, rgb("blue"), 0.24)
    violet_alpha = 0.16 if variant == "violet" else 0.12
    _hero_glow(slide, 0.50, 0.45, 0.34, rgb("violet"), violet_alpha)

    _hero_grid(slide)


def brand_pulse_motif(slide, x=0.18, y=0.52, scale=1.0, color="aqua", transparency=72):
    """Horizontal pulse lines echoing the ParlayVU mark — subtle, not decorative blobs."""
    s = scale
    rails = [
        (x + 0.0, y + 0.08, 0.42 * s, 0.018),
        (x + 0.0, y + 0.20, 0.62 * s, 0.018),
        (x + 0.0, y + 0.32, 0.88 * s, 0.018),
    ]
    for rx, ry, rw, rh in rails:
        rect(slide, rx, ry, rw, rh, "panel2", color, False, 55, 100)
        shape = slide.shapes[-1]
        shape.fill.transparency = transparency
    hub_x, hub_y = x + 0.92 * s, y + 0.20 * s
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(hub_x), Inches(hub_y), Inches(0.11 * s), Inches(0.11 * s))
    dot.fill.solid()
    dot.fill.fore_color.rgb = rgb(color)
    dot.fill.transparency = 35
    dot.line.color.rgb = rgb(color)
    dot.line.transparency = 50
    dot.line.width = Pt(0.75)


def channel_flow_line(slide, x1, y1, x2, y2, color="aqua", width=1.15, transparency=38, dots=True):
    connector(slide, x1, y1, x2, y2, color, width, transparency)
    if dots:
        for t in (0.12, 0.88):
            dx, dy = x1 + (x2 - x1) * t, y1 + (y2 - y1) * t
            mark = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx - 0.045), Inches(dy - 0.045), Inches(0.09), Inches(0.09))
            mark.fill.solid()
            mark.fill.fore_color.rgb = rgb(color)
            mark.fill.transparency = 25
            mark.line.fill.background()


def step_marker(slide, x, y, label, fill="aqua", text_color="ink", size=8.5, d=0.34):
    """Small numbered step badge instead of a large circle."""
    pill(slide, x, y, d, d, label, fill, text_color, size, True, 0)


def metric_block(slide, x, y, w, h, headline, subline=None, fill="aqua", accent="aqua"):
    rect(slide, x, y, w, h, fill, accent, True, 0, 70)
    text = "ink" if fill == "aqua" else "white"
    add_text(slide, x + 0.18, y + 0.14, w - 0.36, 0.42, headline, 22, text, True, font=FONT_HEAD)
    if subline:
        add_text(slide, x + 0.18, y + 0.58, w - 0.36, 0.28, subline, 9.5, "ink" if fill == "aqua" else "muted")


def add_logo(slide, assets, x=0.72, y=0.28, w=1.85):
    if "logo" in assets:
        slide.shapes.add_picture(str(assets["logo"]), Inches(x), Inches(y), width=Inches(w))
    else:
        add_text(slide, x, y + 0.04, w, 0.24, "ParlayVU", 13, "white", True, font=FONT_HEAD)


def add_nav(slide, assets, active="ANGEL"):
    rect(slide, 0.62, 0.18, 11.15, 0.55, "navy", "line", True, 9, 83)
    add_logo(slide, assets)
    nav = [active, "PLATFORM", "PROOF", "PILOTS"]
    x = 8.0
    for item in nav:
        add_text(slide, x, 0.38, 0.75, 0.12, item, 5.7, "muted", True, PP_ALIGN.CENTER)
        x += 0.78


def footer(slide, idx, total):
    add_text(slide, 0.72, 7.05, 4.0, 0.16, "Agentic client-work operating system", 6.8, "soft")
    add_text(slide, 11.8, 7.05, 0.85, 0.16, f"{idx:02d}/{total:02d}", 6.8, "soft", False, PP_ALIGN.RIGHT)


def title_block(slide, kicker, title, subtitle=None):
    add_text(slide, 0.74, 1.14, 4.8, 0.18, kicker.upper(), 6.8, "aqua", True)
    add_text(slide, 0.72, 1.48, 6.4, 1.55, title, 31, "white", True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, 0.76, 3.12, 5.6, 0.58, subtitle, 12.7, "muted")


def card(slide, x, y, w, h, label, body, fill="panel", accent="aqua", body_size=11.2):
    rect(slide, x, y, w, h, fill, "line", True, 5, 78)
    add_text(slide, x + 0.22, y + 0.18, w - 0.44, 0.16, label.upper(), 6.6, accent, True)
    add_text(slide, x + 0.22, y + 0.48, w - 0.44, h - 0.56, body, body_size, "white", False)


def image_card(slide, path: Path, x, y, w, h, title, eyebrow, body):
    rect(slide, x, y, w, h, "panel", "line", True, 0, 72)
    if path and path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h * 0.58))
    add_text(slide, x + 0.22, y + h * 0.62, w - 0.44, 0.15, eyebrow.upper(), 6.4, "aqua", True)
    add_text(slide, x + 0.22, y + h * 0.72, w - 0.44, 0.34, title, 17, "white", True, font=FONT_HEAD)
    add_text(slide, x + 0.24, y + h * 0.86, w - 0.48, 0.42, body, 7.5, "muted", True)


def connector(slide, x1, y1, x2, y2, color="aqua", width=1.4, transparency=45):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.transparency = transparency
    return line


def workflow_hub(slide, x, y, label, accent="aqua"):
    """ParlayVU-style hub: dot + vertical spine for meetings → memory → approvals → outputs."""
    brand_pulse_motif(slide, x - 0.08, y - 0.05, 0.55, accent, 78)
    spine_y = y + 0.42
    for i, item in enumerate(["Meetings", "Memory", "Approvals", "Outputs"]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(spine_y + i * 0.38), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(accent)
        dot.fill.transparency = 20
        dot.line.fill.background()
        if i < 3:
            connector(slide, x + 0.04, spine_y + i * 0.38 + 0.04, x + 0.04, spine_y + (i + 1) * 0.38 + 0.04, accent, 0.9, 58)
        add_text(slide, x + 0.16, spine_y + i * 0.38 - 0.02, 1.1, 0.14, item, 7.2, "muted", True)
    add_text(slide, x - 0.02, y + 1.72, 1.35, 0.14, label.upper(), 6.2, accent, True)


def slide_title(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    brand_pulse_motif(slide, 0.42, 5.55, 0.65, "aqua", 84)
    add_nav(slide, assets)
    title_block(
        slide,
        "Angel / Pre-Seed",
        "AI digital workers for client work.",
        "ParlayVU turns meetings, source material, approvals, and follow-up into an operating system for expert firms.",
    )
    pill(slide, 0.78, 4.05, 1.72, 0.36, "$100k-$300k raise")
    pill(slide, 2.72, 4.05, 2.2, 0.36, "Teams/M365 beachhead", "panel2", "white")
    pill(slide, 5.18, 4.05, 1.8, 0.36, "RamAir proof", "panel2", "white")
    image_card(
        slide,
        assets.get("nathan"),
        7.95,
        1.22,
        3.65,
        3.48,
        "Nathan Ellis",
        "Lead Orchestrator",
        "Coordinates strategy, agent routing, approval gates, and final delivery.",
    )
    add_text(slide, 8.15, 5.65, 3.3, 0.42, "VU is your brand. Parlay your VU.", 17, "white", True, PP_ALIGN.CENTER, FONT_HEAD)
    footer(slide, idx, total)


def slide_problem(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "violet")
    brand_pulse_motif(slide, 11.2, 0.35, 0.9, "violet", 82)
    add_nav(slide, assets, "PROBLEM")
    title_block(slide, "Problem", "Client work leaks margin after every call.", "Generic AI drafts assets. Expert firms still need the whole workflow coordinated.")
    workflow_hub(slide, 5.55, 2.72, "Client call", "violet")
    channel_flow_line(slide, 6.72, 3.18, 7.05, 3.18, "violet", 1.0, 48)
    items = [
        ("Notes", "transcripts, recaps", 1.0, 2.0),
        ("Tasks", "owners, deadlines", 1.3, 4.72),
        ("Files", "source material", 3.7, 1.28),
        ("Approvals", "claims, sends, deploys", 4.05, 5.45),
        ("Reporting", "dashboards, updates", 8.2, 5.2),
        ("Content", "posts, pages, emails", 10.0, 3.12),
        ("Follow-up", "client-ready next steps", 8.62, 1.42),
    ]
    for label, body, x, y in items:
        card(slide, x, y, 2.2, 0.9, label, body, "panel", "blue", 8.7)
        channel_flow_line(slide, 6.35, 3.55, x + 1.1, y + 0.46, "blue", 0.9, 55)
    add_text(slide, 1.25, 6.45, 10.5, 0.28, "The buyer pain is not one prompt. It is memory, routing, governance, and repeatable delivery.", 13.5, "white", True, PP_ALIGN.CENTER, FONT_HEAD)
    footer(slide, idx, total)


def slide_solution(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_nav(slide, assets, "SOLUTION")
    title_block(slide, "Solution", "Nathan runs the client workflow from Teams.", "A front-door orchestrator routes work to specialist agents, grounds it in project memory, and keeps approvals visible.")
    steps = [
        ("Brief Intake", "audience, source, constraints"),
        ("Agent Routing", "right specialist, visible handoff"),
        ("Asset Factory", "notes, sites, emails, dashboards"),
        ("Launch Loop", "approve, deploy, refine"),
    ]
    brand_pulse_motif(slide, 0.55, 2.35, 0.75, "aqua", 80)
    x = 0.82
    for i, (label, body) in enumerate(steps):
        fill = "panel2" if i != 1 else "aqua"
        accent = "aqua" if i != 1 else "ink"
        color = "white" if i != 1 else "ink"
        rect(slide, x, 3.08, 2.78, 1.56, fill, "line", True, 0, 72)
        step_marker(slide, x + 0.18, 3.22, str(i + 1), accent, color, 8.2, 0.3)
        add_text(slide, x + 0.62, 3.32, 2.3, 0.28, label, 17, color, True, font=FONT_HEAD)
        add_text(slide, x + 0.62, 3.95, 2.2, 0.34, body, 9.1, color if i == 1 else "muted")
        if i < len(steps) - 1:
            channel_flow_line(slide, x + 2.86, 3.86, x + 3.04, 3.86, accent, 1.35, 32)
        x += 3.04
    rect(slide, 1.4, 5.48, 10.5, 0.66, "navy", "line", True, 3, 80)
    add_text(slide, 1.75, 5.69, 9.8, 0.18, "Human approval gates govern client-facing sends, publishing, deployments, claims, and commitments.", 11.8, "white", True, PP_ALIGN.CENTER)
    footer(slide, idx, total)


def slide_product(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "violet")
    add_nav(slide, assets, "PRODUCT")
    title_block(slide, "Product", "A Microsoft-native client-work OS.", "Teams and Microsoft 365 are the launch surface customers already use, not the company’s ceiling.")
    layers = [
        ("Front doors", "Teams messages | M365 Files | Word | email/task surfaces | web/demo UI"),
        ("Orchestration", "Nathan | FastAPI/LangGraph | specialist agent registry | readiness checks"),
        ("Memory + governance", "clients | projects | source assets | approval states | agent events"),
        ("Output systems", "meeting notes | approval packets | campaign sites | dashboards | follow-up"),
    ]
    y = 2.56
    for label, body in layers:
        rect(slide, 1.0, y, 11.1, 0.75, "panel", "line", True, 5, 76)
        add_text(slide, 1.28, y + 0.2, 2.0, 0.15, label.upper(), 6.5, "aqua", True)
        add_text(slide, 3.15, y + 0.18, 8.5, 0.2, body, 10.4, "white")
        y += 0.9
    pill(slide, 1.0, 6.25, 5.55, 0.36, "Working proof: routing, memory scaffold, Files outputs, approvals", "aqua", "ink", 8.2)
    pill(slide, 6.82, 6.25, 5.3, 0.36, "Roadmap: native Teams audio/video media bridge", "panel2", "white", 8.2)
    footer(slide, idx, total)


def slide_exists(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_nav(slide, assets, "PROOF")
    title_block(slide, "What Exists Now", "The demo proves the operating loop.", "Current proof is specific: workflow infrastructure, seeded pilot artifacts, and approval-aware outputs.")
    proof = [
        ("Nathan routing", "Teams-style messages against project memory"),
        ("M365 Files", "meeting notes as .md and .docx"),
        ("Approvals", "endpoints and Teams-ready cards"),
        ("RamAir workspace", "channel model, artifacts, dashboard starter"),
        ("Avatar validation", "provider-hosted Tavus path; native Teams A/V remains roadmap"),
        ("Readiness checks", "LLM, memory, M365, Teams, avatar, approvals"),
    ]
    x, y = 0.92, 2.48
    for i, (label, body) in enumerate(proof):
        card(slide, x, y, 3.75, 1.02, label, body, "panel", "aqua" if i < 4 else "blue", 9.4)
        x += 4.03
        if x > 9:
            x, y = 0.92, y + 1.28
    add_text(slide, 1.12, 6.25, 10.9, 0.32, "Credibility boundary: do not imply mature SaaS traction or production native Teams media until pilots and media bridge validation prove it.", 11.4, "muted", False, PP_ALIGN.CENTER)
    footer(slide, idx, total)


def slide_ramair(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_nav(slide, assets, "CASE STUDY")
    title_block(slide, "ParlayVU In The Field", "RamAir turns one weekly podcast into a content engine.", "The first client story shows the repeatable pattern: source material becomes governed outputs across channels.")
    image_card(
        slide,
        assets.get("ramair"),
        0.92,
        2.38,
        3.1,
        3.55,
        "Straight from the Hart",
        "RamAir / ramair.co",
        "Podcast source material becomes a coordinated workspace and distribution cadence.",
    )
    outputs = [
        ("YouTube + site", 4.55, 2.28),
        ("Short clips", 6.92, 2.28),
        ("Social posts", 9.28, 2.28),
        ("Transcript insights", 4.55, 3.78),
        ("Case studies", 6.92, 3.78),
        ("Sell sheets", 9.28, 3.78),
    ]
    for label, x, y in outputs:
        channel_flow_line(slide, 4.02, 3.85, x + 0.98, y + 0.48, "aqua", 0.95, 52)
        card(slide, x, y, 1.95, 0.96, "Output", label, "panel", "aqua", 9.5)
    rect(slide, 4.65, 5.42, 6.7, 0.72, "navy", "line", True, 4, 80)
    add_text(slide, 4.95, 5.63, 6.1, 0.18, "Use as pilot proof of a workflow pattern, not broad SaaS traction.", 10.4, "muted", True, PP_ALIGN.CENTER)
    footer(slide, idx, total)


def slide_market_model(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "violet")
    add_nav(slide, assets, "MARKET")
    title_block(slide, "Market + Model", "Start where expertise becomes deliverables.", "Founder-led pilots reveal the repeatable workspace subscription modules.")
    card(slide, 0.95, 2.44, 3.5, 2.62, "Ideal early customer", "Boutique agencies, consultants, coaches, professional services firms, and expert-led businesses already using Microsoft 365.", "panel", "aqua", 12)
    card(slide, 4.92, 2.44, 3.5, 2.62, "Initial promise", "Turn client meetings and source material into approved deliverables, updates, and follow-up inside the tools the firm already uses.", "panel", "blue", 12)
    card(slide, 8.88, 2.44, 3.5, 2.62, "Revenue path", "Pilot setup fee plus monthly subscription by clients/projects, agent seats, integrations, outputs, and premium workflow modules.", "panel", "violet", 12)
    add_text(slide, 1.15, 6.04, 10.95, 0.36, "Ann Arbor/Michigan is the launch network and credibility base, not the market ceiling.", 13, "white", True, PP_ALIGN.CENTER, FONT_HEAD)
    footer(slide, idx, total)


def slide_gtm(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_nav(slide, assets, "GTM")
    title_block(slide, "Go To Market", "Founder-led pilots create the first playbook.", "Sell one concrete workflow, measure the outputs, then template the repeatable pieces.")
    steps = [
        ("Introductions", "angels, operators, agencies"),
        ("RamAir-style pilot", "one workflow, visible outputs"),
        ("Proof metrics", "time saved, output cadence"),
        ("Templates", "notes, packets, dashboards"),
        ("Referrals", "adjacent firms, same pain"),
    ]
    x = 0.72
    for i, (label, body) in enumerate(steps):
        step_marker(slide, x + 0.92, 2.58, str(i + 1), "aqua" if i == 0 else "panel2", "ink" if i == 0 else "white", 8.0, 0.3)
        card(slide, x, 3.05, 2.25, 1.08, label, body, "panel", "aqua", 8.6)
        if i < len(steps) - 1:
            channel_flow_line(slide, x + 2.28, 3.58, x + 2.52, 3.58, "aqua", 1.1, 40)
        x += 2.52
    add_text(slide, 1.08, 5.95, 11.1, 0.4, "3-5 pilots are enough to convert a services-assisted motion into a productized client-workspace package.", 14, "white", True, PP_ALIGN.CENTER, FONT_HEAD)
    footer(slide, idx, total)


def slide_differentiation(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_nav(slide, assets, "DIFFERENTIATION")
    title_block(slide, "Differentiation", "Not another blank chat box.", "ParlayVU combines memory, governance, and outputs in the Microsoft workflow where client work already happens.")
    card(slide, 1.0, 2.55, 4.65, 1.18, "Generic AI", "drafts content, but loses workflow context", "panel", "blue", 11)
    card(slide, 1.0, 4.28, 4.65, 1.18, "Content tools", "help make assets, not client operations", "panel", "blue", 11)
    rect(slide, 7.05, 2.55, 4.95, 2.9, "aqua", "aqua", True, 0, 95)
    add_text(slide, 7.45, 2.9, 4.2, 0.24, "PARLAYVU", 7, "ink", True)
    add_text(slide, 7.42, 3.32, 4.1, 0.82, "Coordinates client work across memory, agents, approvals, and outputs.", 20, "ink", True, font=FONT_HEAD)
    add_text(slide, 7.44, 4.72, 4.0, 0.34, "Teams/M365 is the wedge. The durable value is governed agentic workflow.", 9.5, "ink", True)
    connector(slide, 6.32, 2.35, 6.32, 5.82, "line", 1.6, 74)
    footer(slide, idx, total)


def slide_roadmap(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "violet")
    add_nav(slide, assets, "ROADMAP")
    title_block(slide, "Roadmap", "90 days to pilot-ready product.", "Sequence core workflow hardening before making bigger native media claims.")
    milestones = [
        ("30 days", "Polish RamAir workspace, notes, approvals, readiness, demo script"),
        ("60 days", "Package onboarding templates and run 3-5 qualified pilot conversations"),
        ("90 days", "Paid pilots or LOIs, measurable outputs, pricing and module clarity"),
        ("Later", "Validate native Teams roster/media bridge after Graph and provider terms"),
    ]
    y_line = 3.18
    connector(slide, 1.35, y_line, 11.55, y_line, "aqua", 1.25, 48)
    x = 0.95
    for i, (label, body) in enumerate(milestones):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.62), Inches(y_line - 0.04), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb("aqua" if label != "Later" else "blue")
        dot.fill.transparency = 18
        dot.line.fill.background()
        card(slide, x, 2.72, 2.62, 1.32, label, body, "panel", "aqua" if label != "Later" else "blue", 8.9)
        x += 3.0
    footer(slide, idx, total)


def slide_ask(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_nav(slide, assets, "ASK")
    title_block(slide, "Ask", "Raising $100k-$300k to turn proof into revenue.", "Capital plus pilot/customer introductions are equally valuable at this stage.")
    metric_block(slide, 0.88, 2.38, 2.85, 1.35, "$100k–$300k", "Angel / pre-seed raise")
    brand_pulse_motif(slide, 11.05, 5.1, 0.7, "aqua", 85)
    uses = [
        ("Product hardening", 4.05, 2.28),
        ("Pilot delivery", 6.92, 2.28),
        ("M365 / Teams integration", 9.78, 2.28),
        ("Avatar validation", 4.05, 4.05),
        ("Design + pitch assets", 6.92, 4.05),
        ("Founder-led GTM", 9.78, 4.05),
    ]
    for label, x, y in uses:
        pill(slide, x, y, 2.34, 0.5, label, "panel2", "white", 8.6)
    rect(slide, 1.1, 5.82, 11.0, 0.62, "navy", "line", True, 4, 80)
    add_text(slide, 1.4, 6.02, 10.38, 0.18, "Target: RamAir polished, 3-5 qualified pilots active, first paid pilots or signed LOIs within 90 days of funding.", 10.7, "white", True, PP_ALIGN.CENTER)
    footer(slide, idx, total)


def slide_close(prs, assets, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_nav(slide, assets, "START")
    add_text(slide, 1.02, 1.35, 8.6, 1.16, "VU is your brand.\nParlay your VU.", 40, "white", True, font=FONT_HEAD)
    add_text(slide, 1.08, 3.25, 5.85, 0.7, "ParlayVU turns every client interaction into approved work, reusable memory, and distribution-ready output.", 17.2, "muted")
    rect(slide, 7.7, 2.2, 3.9, 2.45, "panel", "line", True, 2, 75)
    add_text(slide, 8.08, 2.65, 3.12, 0.32, "Investor asks", 21, "white", True, PP_ALIGN.CENTER, FONT_HEAD)
    add_text(slide, 8.25, 3.35, 2.8, 0.75, "Small checks\nPilot introductions\nMicrosoft 365 operator/customer intros", 13.2, "muted", False, PP_ALIGN.CENTER)
    pill(slide, 8.44, 4.98, 2.38, 0.38, "Book a strategy call")
    footer(slide, idx, total)


BUILDERS = [
    slide_title,
    slide_problem,
    slide_solution,
    slide_product,
    slide_exists,
    slide_ramair,
    slide_market_model,
    slide_gtm,
    slide_differentiation,
    slide_roadmap,
    slide_ask,
    slide_close,
]


def validate(path: Path, expected: int) -> None:
    with ZipFile(path) as zf:
        names = zf.namelist()
        slide_xml = [name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        if len(slide_xml) != expected:
            raise RuntimeError(f"Expected {expected} slide XML files, found {len(slide_xml)}")
        if "ppt/presentation.xml" not in names:
            raise RuntimeError("Missing presentation.xml")
    prs = Presentation(str(path))
    if len(prs.slides) != expected:
        raise RuntimeError(f"python-pptx reload found {len(prs.slides)} slides")


def build() -> Path:
    assets = ensure_assets()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    total = len(BUILDERS)
    for i, builder in enumerate(BUILDERS, start=1):
        builder(prs, assets, i, total)
    prs.save(OUTPUT)
    validate(OUTPUT, total)
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out}")
    print(f"Slides: {len(BUILDERS)}")
