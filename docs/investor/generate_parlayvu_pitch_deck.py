from __future__ import annotations

from pathlib import Path
from urllib.request import Request, urlopen
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "assets"
OUTPUT = OUT_DIR / "ParlayVU_Angel_Pitch_Deck.pptx"
LOGO_PATH = ASSET_DIR / "parlayvu-logo-dark.png"
LOGO_URL = "https://parlayvu.ai/brand/parlayvu-logo-dark.png"


BRAND = {
    "deep_green": RGBColor(32, 53, 47),      # #20352f
    "green_hover": RGBColor(48, 77, 69),     # #304d45
    "muted": RGBColor(82, 100, 93),          # #52645d
    "copper": RGBColor(162, 95, 43),         # #a25f2b
    "gold": RGBColor(217, 144, 79),          # #d9904f
    "sand": RGBColor(243, 215, 166),         # #f3d7a6
    "sage": RGBColor(127, 155, 143),         # #7f9b8f
    "cream": RGBColor(249, 245, 237),
    "white": RGBColor(255, 255, 255),
    "dark": RGBColor(22, 32, 29),
}

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


SLIDES = [
    {
        "eyebrow": "Angel / Pre-Seed Pitch",
        "title": "ParlayVU.ai turns client work into an AI-operated workspace",
        "subtitle": "Microsoft-native AI operating system for client-facing expert firms.",
        "bullets": [
            "Nathan orchestrates specialist AI agents inside Teams, Microsoft 365, and project memory.",
            "Raising $100k-$300k to convert a working product into repeatable pilots.",
        ],
        "cue": "Open with practical B2B productivity: firms already live in Microsoft 365; ParlayVU makes that work surface agentic.",
        "layout": "title",
    },
    {
        "eyebrow": "Problem",
        "title": "Client-facing firms lose margin after every call",
        "bullets": [
            "Calls, transcripts, notes, assets, tasks, dashboards, emails, and approvals scatter across disconnected tools.",
            "Generic AI can draft content, but it does not remember the client or manage the workflow.",
            "Boutique agencies, consultants, coaches, and expert businesses need AI leverage without building an internal AI team.",
        ],
        "cue": "Emphasize that the workflow is the product for expert firms; fragmentation is a margin leak, not just an annoyance.",
    },
    {
        "eyebrow": "Solution",
        "title": "Nathan runs the client workflow from Teams",
        "bullets": [
            "Nathan is the Teams front door and orchestrator for a specialist agent team.",
            "Agents turn source material into notes, campaign assets, websites, approval packets, dashboards, and follow-up.",
            "Human approvals remain the gate for client-facing sends, publishing, deployments, claims, and commitments.",
        ],
        "cue": "Position ParlayVU as governed execution, not an uncontrolled autonomous employee.",
    },
    {
        "eyebrow": "Product",
        "title": "A Microsoft-native AI workspace, not another chatbot",
        "bullets": [
            "Teams channels bind to client and project memory, so Nathan answers in context.",
            "Teams Files publishing creates meeting notes as .md for memory and .docx for client review.",
            "Graph, agent mailboxes, approval cards, Planner direction, SharePoint Files, and Power BI starter data form the operating surface.",
            "Dylan can generate Astro/Tailwind campaign sites and route deployment through approvals.",
        ],
        "cue": "Make the Microsoft-native point concrete: Teams, SharePoint Files, Word, Planner direction, Power BI starters, and approvals.",
    },
    {
        "eyebrow": "Current Proof",
        "title": "The demo already proves the workflow loop",
        "bullets": [
            "FastAPI/LangGraph backend with readiness checks for LLM, memory, M365, Teams, avatar, and approval configuration.",
            "Teams-style Nathan message routing works against project memory.",
            "RamAir pilot workspace documents channels, Files folders, meeting notes, approvals, and dashboard starter data.",
            "Tavus validates a provider-hosted Nathan conversation path; native Teams audio/video remains roadmap/scaffold.",
        ],
        "cue": "Be crisp about proof boundaries. The working loop exists; production native Teams avatar participation is not claimed.",
        "layout": "proof",
    },
    {
        "eyebrow": "Why Now",
        "title": "AI agents are moving into the tools where work already happens",
        "bullets": [
            "Microsoft Teams and Microsoft 365 are the default operating surface for many client-facing firms.",
            "AI has moved from single prompts to agentic workflows, but most SMB teams cannot assemble the plumbing themselves.",
            "Buyers want productivity gains with control: project memory, approvals, auditability, and client-ready outputs.",
        ],
        "cue": "Tie the timing to practical adoption: teams want leverage, but they need trust and workflow fit.",
    },
    {
        "eyebrow": "Market Wedge",
        "title": "Start with firms where expertise becomes client deliverables",
        "bullets": [
            "Initial wedge: boutique agencies, consultants, coaches, professional services firms, and expert-led businesses.",
            "Ideal early customer: Microsoft 365/Teams user with repeatable meetings, content, campaigns, reporting, and approvals.",
            "Ann Arbor/Michigan provides a practical B2B launch network without limiting the broader market.",
        ],
        "cue": "Show discipline: the first wedge is narrow enough for founder-led pilots and broad enough to become a repeatable category.",
    },
    {
        "eyebrow": "Business Model",
        "title": "High-touch pilots now, repeatable workspace subscriptions next",
        "bullets": [
            "Early revenue: setup or pilot fee plus monthly subscription.",
            "Subscription dimensions: active clients/projects, agent seats, integrations, generated outputs, storage, and premium modules.",
            "Services-assisted onboarding exposes repeatable templates and accelerates learning.",
        ],
        "cue": "Angels should hear that services are a learning wedge, not the long-term ceiling.",
    },
    {
        "eyebrow": "Go To Market",
        "title": "Founder-led pilots create the first repeatable playbook",
        "bullets": [
            "Run 3-5 RamAir-style pilots with firms already using Microsoft 365.",
            "Sell around one concrete workflow: turn client meetings and source material into approved deliverables and follow-up.",
            "Convert successful pilots into templates: channel setup, meeting notes, approval packets, campaign kits, dashboards, and weekly updates.",
        ],
        "cue": "Ask investors for customer introductions as much as capital: agencies, consultants, Microsoft-centric SMBs, and operators.",
    },
    {
        "eyebrow": "Differentiation",
        "title": "Project memory, approvals, and Microsoft 365 outputs compound together",
        "bullets": [
            "Not generic chat: Teams-native project memory and client workspace execution.",
            "Not a content-only tool: agents coordinate notes, sites, drafts, approvals, dashboards, and follow-up.",
            "Not uncontrolled autonomy: approval gates govern client-facing actions.",
            "Roadmap advantage: provider-hosted avatar proof plus conservative native Teams media-bot scaffold.",
        ],
        "cue": "Contrast against generic AI and content tools by grounding the moat in memory, governance, and Microsoft-native outputs.",
        "layout": "three_cards",
    },
    {
        "eyebrow": "Roadmap",
        "title": "The next 90 days turn the demo into pilot-ready product",
        "bullets": [
            "30 days: polish RamAir pilot workspace, meeting-note flow, approval cards, readiness checks, and demo script.",
            "60 days: package onboarding templates for 3-5 pilot firms and tighten M365/Teams Files setup.",
            "90 days: convert pilot learnings into workspace modules, dashboard updates, stronger memory ingestion, and early pricing.",
            "Later: validate native Teams roster participation and audio/video bridge after Graph media and provider terms are proven.",
        ],
        "cue": "Use the roadmap to reinforce credibility: the media bridge is future validation, not current production capability.",
        "layout": "timeline",
    },
    {
        "eyebrow": "Ask",
        "title": "Raising $100k-$300k to turn working product into early revenue",
        "bullets": [
            "Use of funds: product hardening, pilot delivery, M365/Teams integration, Tavus/LiveAvatar validation, design assets, and founder-led GTM.",
            "By close: RamAir polished, 3-5 qualified pilot conversations active, and repeatable workspace templates documented.",
            "Within 90 days of funding: first paid pilots or signed pilot LOIs, measurable workflow outputs, and clearer subscription packaging.",
            "Investor asks: small checks, pilot introductions, Microsoft 365 operator/customer intros, and Michigan angel referrals.",
        ],
        "cue": "Close on the conversion: this round turns a working demo and pilot scaffold into repeatable customer pilots.",
        "layout": "ask",
    },
]


def add_textbox(slide, x, y, w, h, text="", font_size=18, color=None, bold=False, font=FONT_BODY, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.name = font
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color or BRAND["deep_green"]
    return box


def add_footer(slide, slide_no, total, logo_available):
    add_textbox(slide, 0.55, 7.08, 2.6, 0.18, "ParlayVU.ai", 7, BRAND["muted"], True)
    add_textbox(slide, 11.6, 7.08, 0.8, 0.18, f"{slide_no:02d}/{total:02d}", 7, BRAND["muted"], False, align=PP_ALIGN.RIGHT)
    if not logo_available and slide_no == 1:
        add_textbox(slide, 8.55, 7.05, 3.0, 0.2, "Logo fallback: editable text treatment", 6.5, BRAND["muted"], False, align=PP_ALIGN.RIGHT)


def add_brand_mark(slide, logo_available):
    if logo_available:
        try:
            slide.shapes.add_picture(str(LOGO_PATH), Inches(0.6), Inches(0.34), width=Inches(1.6))
            return
        except Exception:
            pass
    add_textbox(slide, 0.62, 0.36, 1.8, 0.28, "ParlayVU", 13, BRAND["deep_green"], True, FONT_HEAD)


def add_background(slide, variant="light"):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BRAND["cream"] if variant == "light" else BRAND["deep_green"]
    if variant == "light":
        orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(-0.55), Inches(3.8), Inches(3.8))
        orb.fill.solid()
        orb.fill.fore_color.rgb = BRAND["sand"]
        orb.line.fill.background()
        orb.fill.transparency = 25
        ribbon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(0.72), Inches(2.4), Inches(0.26))
        ribbon.fill.solid()
        ribbon.fill.fore_color.rgb = BRAND["gold"]
        ribbon.line.fill.background()
        ribbon.rotation = -14
        ribbon.fill.transparency = 6


def add_title(slide, eyebrow, title, subtitle=None):
    add_textbox(slide, 0.85, 0.92, 6.5, 0.28, eyebrow.upper(), 9.5, BRAND["copper"], True)
    add_textbox(slide, 0.83, 1.26, 7.85, 1.35, title, 28, BRAND["deep_green"], True, FONT_HEAD)
    if subtitle:
        add_textbox(slide, 0.88, 2.48, 6.55, 0.55, subtitle, 15, BRAND["muted"])


def add_bullets(slide, bullets, x=0.92, y=3.05, w=6.55, h=2.7, font_size=13.6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = BRAND["dark"]
        p.space_after = Pt(8)
    return box


def add_cue(slide, cue):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(6.02), Inches(11.72), Inches(0.68))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND["white"]
    shape.fill.transparency = 4
    shape.line.color.rgb = BRAND["sand"]
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    p = tf.paragraphs[0]
    p.text = "Founder cue: " + cue
    p.font.name = FONT_BODY
    p.font.size = Pt(8.5)
    p.font.color.rgb = BRAND["muted"]


def add_side_panel(slide, headline, items):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.42), Inches(1.23), Inches(3.55), Inches(4.53))
    panel.fill.solid()
    panel.fill.fore_color.rgb = BRAND["deep_green"]
    panel.line.fill.background()
    add_textbox(slide, 8.75, 1.62, 2.75, 0.6, headline, 17, BRAND["white"], True, FONT_HEAD)
    y = 2.48
    for label, detail in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.82), Inches(y + 0.05), Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BRAND["gold"]
        dot.line.fill.background()
        add_textbox(slide, 9.08, y - 0.03, 2.45, 0.23, label, 8.5, BRAND["sand"], True)
        add_textbox(slide, 9.08, y + 0.22, 2.42, 0.42, detail, 8.4, RGBColor(230, 238, 232))
        y += 0.95


def add_metric_card(slide, x, y, label, value, detail):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.08), Inches(1.08))
    card.fill.solid()
    card.fill.fore_color.rgb = BRAND["white"]
    card.line.color.rgb = BRAND["sand"]
    add_textbox(slide, x + 0.18, y + 0.14, 1.65, 0.24, label.upper(), 6.8, BRAND["copper"], True)
    add_textbox(slide, x + 0.18, y + 0.4, 1.7, 0.34, value, 17, BRAND["deep_green"], True, FONT_HEAD)
    add_textbox(slide, x + 0.18, y + 0.76, 1.62, 0.18, detail, 6.7, BRAND["muted"])


def add_standard_visual(slide, index, layout):
    if layout == "proof":
        add_side_panel(slide, "Current proof, stated carefully", [
            ("Working", "Nathan routing, memory scaffold, M365 Files outputs, approvals"),
            ("Demo-ready", "RamAir workspace, readiness checks, approval workflow"),
            ("Validated", "Provider-hosted avatar path through Tavus"),
            ("Roadmap", "Native Teams media bridge and production A/V participation"),
        ])
        return
    if layout == "three_cards":
        labels = [
            ("Memory", "Client/project context lives with the workflow."),
            ("Governance", "Approval gates protect client-facing actions."),
            ("Outputs", "Work lands in Teams, Files, dashboards, sites, and follow-up."),
        ]
        y = 1.52
        for label, detail in labels:
            add_metric_card(slide, 8.52, y, "Differentiator", label, detail)
            y += 1.28
        return
    if layout == "timeline":
        add_side_panel(slide, "Milestone cadence", [
            ("30 days", "Polish RamAir and demo loop"),
            ("60 days", "Package pilot onboarding templates"),
            ("90 days", "Paid pilots or LOIs plus pricing clarity"),
            ("Later", "Validate native Teams media path"),
        ])
        return
    if layout == "ask":
        add_side_panel(slide, "Round shape", [
            ("Raise", "$100k-$300k angel / pre-seed"),
            ("Capital", "Product hardening and pilot delivery"),
            ("Introductions", "Pilot customers and Microsoft-centric SMBs"),
            ("Milestone", "Working demo to first revenue evidence"),
        ])
        return

    visuals = [
        ("Surface", "Teams + Microsoft 365"),
        ("Control", "Approval-gated workflow"),
        ("Output", "Client-ready deliverables"),
    ]
    add_side_panel(slide, "Operating loop", visuals)


def build_deck():
    logo_available = fetch_logo()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    total = len(SLIDES)
    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        add_background(slide, "light")
        add_brand_mark(slide, logo_available)

        if data.get("layout") == "title":
            add_title(slide, data["eyebrow"], data["title"], data.get("subtitle"))
            add_bullets(slide, data["bullets"], 0.92, 3.35, 6.6, 1.35, 14.2)
            add_side_panel(slide, "Brand cues from parlayvu.ai", [
                ("Palette", "Deep green, copper, sand, sage"),
                ("Style", "Rounded cards, editorial discipline"),
                ("Type", "Clean sans-serif PPT-native Aptos"),
                ("Logo", "Downloaded if public asset is accessible"),
            ])
        else:
            add_title(slide, data["eyebrow"], data["title"])
            add_bullets(slide, data["bullets"])
            add_standard_visual(slide, idx, data.get("layout"))

        add_cue(slide, data["cue"])
        add_footer(slide, idx, total, logo_available)

    prs.save(OUTPUT)
    validate_pptx(OUTPUT, total)
    return logo_available


def fetch_logo() -> bool:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    if LOGO_PATH.exists() and LOGO_PATH.stat().st_size > 0:
        return True
    try:
        request = Request(LOGO_URL, headers={"User-Agent": "Mozilla/5.0"})
        with urlopen(request, timeout=15) as response:
            data = response.read()
        if data.startswith(b"\x89PNG"):
            LOGO_PATH.write_bytes(data)
            return True
    except Exception as exc:
        (OUT_DIR / "ParlayVU_Angel_Pitch_Deck_asset_notes.txt").write_text(
            "Could not fetch public logo asset from "
            f"{LOGO_URL}. Used editable text-logo fallback in deck.\nError: {exc}\n",
            encoding="utf-8",
        )
    return False


def validate_pptx(path: Path, expected_slides: int):
    with ZipFile(path) as zf:
        slide_files = [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
    if len(slide_files) != expected_slides:
        raise RuntimeError(f"Expected {expected_slides} slides, found {len(slide_files)}")

    reloaded = Presentation(str(path))
    if len(reloaded.slides) != expected_slides:
        raise RuntimeError(f"python-pptx reload found {len(reloaded.slides)} slides")


if __name__ == "__main__":
    used_logo = build_deck()
    print(f"Wrote {OUTPUT}")
    print(f"Slides: {len(SLIDES)}")
    print(f"Logo asset: {'downloaded' if used_logo else 'text fallback'}")
